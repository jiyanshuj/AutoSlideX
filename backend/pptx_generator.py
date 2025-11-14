from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
from typing import Dict, List
from PIL import Image
import os
import re

class PPTXGenerator:
    """
    Modern PowerPoint generator with split-screen image + content layout
    Enhanced with smart image detection using Pixabay API
    """
    
    def __init__(self, template="modern"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.template = template
        self.theme_colors = self._get_theme_colors(template)
        
        # ✅ FIX: Initialize ALL API keys
        self.google_api_key = os.getenv("GOOGLE_API_KEY", "")
        self.google_cx = os.getenv("GOOGLE_SEARCH_ENGINE_ID", "")
        self.pixabay_api_key = os.getenv("PIXABAY_API_KEY", "")  # ✅ ADDED
        
        # Diagram and technical content detection patterns
        self.diagram_patterns = {
            # UML Diagrams
            r'\b(class\s+diagram|uml\s+class)\b': 'UML class diagram computer',
            r'\b(sequence\s+diagram|uml\s+sequence)\b': 'UML sequence diagram software',
            r'\b(use[\s-]?case\s+diagram|uml\s+use[\s-]?case)\b': 'UML use case diagram software',
            r'\b(activity\s+diagram|uml\s+activity)\b': 'UML activity diagram workflow',
            r'\b(state\s+diagram|state\s+machine)\b': 'state machine diagram software',
            r'\b(component\s+diagram|uml\s+component)\b': 'software component diagram',
            r'\b(deployment\s+diagram|uml\s+deployment)\b': 'deployment diagram infrastructure',
            
            # Operating System Specific
            r'\b(layered\s+(architecture|structure|approach|design).*os|os.*layered)\b': 'operating system layers diagram illustration',
            r'\b(os\s+architecture|operating\s+system\s+architecture)\b': 'operating system architecture layers diagram',
            r'\b(kernel\s+architecture|kernel\s+layers)\b': 'kernel architecture diagram layers',
            r'\b(os\s+structure|operating\s+system\s+structure)\b': 'operating system structure layers diagram',
            
            # Architecture & System Diagrams (Software focus)
            r'\b(software\s+architecture|application\s+architecture)\b': 'software architecture diagram illustration',
            r'\b(system\s+architecture)\b': 'system architecture diagram technology',
            r'\b(microservices?\s+architecture)\b': 'microservices architecture diagram',
            r'\b(cloud\s+architecture)\b': 'cloud computing architecture diagram',
            r'\b(network\s+architecture)\b': 'network topology diagram',
            r'\b(layered\s+architecture|tier\s+architecture|n-tier)\b': 'layered architecture diagram software illustration',
            
            # Database Diagrams
            r'\b(er\s+diagram|entity[\s-]?relationship)\b': 'ER entity relationship database diagram',
            r'\b(database\s+schema|db\s+schema)\b': 'database schema diagram illustration',
            r'\b(data\s+model|database\s+model)\b': 'database model diagram',
            
            # Network Diagrams
            r'\b(network\s+diagram|network\s+topology)\b': 'network topology diagram illustration',
            r'\b(network\s+infrastructure)\b': 'network infrastructure diagram',
            
            # Process & Flow Diagrams
            r'\b(flowchart|flow\s+chart)\b': 'flowchart process diagram illustration',
            r'\b(process\s+flow|workflow)\b': 'workflow process diagram',
            r'\b(business\s+process)\b': 'business process flow diagram',
            r'\b(data\s+flow|dfd)\b': 'data flow diagram software',
            
            # Project Management
            r'\b(gantt\s+chart)\b': 'gantt chart project timeline',
            r'\b(timeline|project\s+timeline)\b': 'project timeline diagram',
            r'\b(roadmap|product\s+roadmap)\b': 'product roadmap diagram',
            
            # Data Visualization
            r'\b(pie\s+chart)\b': 'pie chart visualization',
            r'\b(bar\s+chart|bar\s+graph)\b': 'bar chart visualization',
            r'\b(line\s+chart|line\s+graph)\b': 'line chart visualization',
            r'\b(scatter\s+plot)\b': 'scatter plot visualization',
            r'\b(histogram)\b': 'histogram chart',
            
            # Technical Diagrams
            r'\b(circuit\s+diagram)\b': 'electronic circuit diagram',
            r'\b(wiring\s+diagram)\b': 'wiring diagram',
            r'\b(block\s+diagram)\b': 'block diagram technical',
            r'\b(organizational\s+chart|org\s+chart)\b': 'organizational chart structure',
            r'\b(mind\s+map)\b': 'mind map diagram',
            r'\b(tree\s+diagram|hierarchy)\b': 'tree hierarchy diagram',
            r'\b(venn\s+diagram)\b': 'venn diagram',
            
            # Software Development
            r'\b(api\s+design|api\s+architecture)\b': 'API architecture diagram',
            r'\b(software\s+stack|tech\s+stack)\b': 'technology stack diagram',
            r'\b(ci[\s/]?cd\s+pipeline)\b': 'CI CD pipeline diagram',
            r'\b(git\s+workflow|git\s+flow)\b': 'git workflow diagram',
            
            # Infrastructure
            r'\b(server\s+architecture)\b': 'server architecture diagram',
            r'\b(container\s+architecture|docker)\b': 'docker container architecture',
            r'\b(kubernetes|k8s)\b': 'kubernetes architecture diagram',
            
            # Operating Systems
            r'\b(operating\s+system|os\s+architecture)\b': 'operating system architecture',
            r'\b(kernel\s+architecture)\b': 'kernel architecture diagram',
            r'\b(memory\s+management)\b': 'memory management diagram',
            r'\b(process\s+scheduling)\b': 'process scheduling diagram',
        }
        
        # Technical topic patterns (for fallback when no diagram is detected)
        self.tech_topics = {
            r'\b(python|java|javascript|c\+\+|ruby|php)\b': 'programming code',
            r'\b(machine\s+learning|deep\s+learning|ai|neural\s+network)\b': 'artificial intelligence',
            r'\b(data\s+science|analytics|big\s+data)\b': 'data science',
            r'\b(web\s+development|frontend|backend)\b': 'web development',
            r'\b(cloud\s+computing|aws|azure|gcp)\b': 'cloud computing',
            r'\b(cybersecurity|security|encryption)\b': 'cybersecurity',
            r'\b(blockchain|cryptocurrency)\b': 'blockchain technology',
            r'\b(iot|internet\s+of\s+things)\b': 'IoT technology',
            r'\b(devops|automation)\b': 'DevOps',
            r'\b(mobile\s+development|android|ios)\b': 'mobile development',
        }
    
    def _get_theme_colors(self, template):
        """Get color scheme based on template"""
        themes = {
            "modern": {
                "primary": RGBColor(41, 128, 185),
                "secondary": RGBColor(52, 73, 94),
                "accent": RGBColor(231, 76, 60),
                "background": RGBColor(248, 249, 250),
                "text": RGBColor(33, 33, 33),
                "light_text": RGBColor(255, 255, 255),
                "dark_text": RGBColor(33, 33, 33)
            },
            "professional": {
                "primary": RGBColor(31, 58, 147),
                "secondary": RGBColor(67, 97, 238),
                "accent": RGBColor(76, 175, 80),
                "background": RGBColor(255, 255, 255),
                "text": RGBColor(33, 33, 33),
                "light_text": RGBColor(255, 255, 255),
                "dark_text": RGBColor(33, 33, 33)
            },
            "creative": {
                "primary": RGBColor(156, 39, 176),
                "secondary": RGBColor(233, 30, 99),
                "accent": RGBColor(255, 193, 7),
                "background": RGBColor(255, 255, 255),
                "text": RGBColor(66, 66, 66),
                "light_text": RGBColor(255, 255, 255),
                "dark_text": RGBColor(33, 33, 33)
            }
        }
        return themes.get(template, themes["modern"])
    
    def _detect_diagram_type(self, text: str) -> str:
        """
        Detect diagram type from text content using pattern matching
        Returns optimized search query for Pixabay
        """
        if not text:
            return None
        
        text_lower = text.lower()
        
        # Check for diagram patterns (priority)
        for pattern, query in self.diagram_patterns.items():
            if re.search(pattern, text_lower, re.IGNORECASE):
                print(f"🎯 Detected diagram type: {query}")
                return query
        
        return None
    
    def _detect_tech_topic(self, text: str) -> str:
        """
        Detect technical topic from text content
        Returns relevant search term
        """
        if not text:
            return None
        
        text_lower = text.lower()
        
        # Check for technical topics
        for pattern, topic in self.tech_topics.items():
            if re.search(pattern, text_lower, re.IGNORECASE):
                print(f"🔍 Detected tech topic: {topic}")
                return topic
        
        return None
    
    def _build_smart_query(self, title: str, content: List[str], provided_query: str = None) -> str:
        """
        Build intelligent search query based on content analysis
        Priority: provided_query > diagram detection > tech topic > title
        """
        # If user provided specific query, use it
        if provided_query and provided_query.strip():
            print(f"📌 Using provided query: {provided_query}")
            return provided_query.strip()
        
        # Combine title and content for analysis
        full_text = f"{title} {' '.join(content)}"
        
        # Step 1: Check for diagram types (highest priority)
        diagram_query = self._detect_diagram_type(full_text)
        if diagram_query:
            return diagram_query
        
        # Step 2: Check for technical topics
        tech_query = self._detect_tech_topic(full_text)
        if tech_query:
            return tech_query
        
        # Step 3: Fall back to title-based query
        # Extract meaningful keywords from title
        keywords = title.lower().split()
        meaningful_keywords = [w for w in keywords if len(w) > 3][:3]
        
        if meaningful_keywords:
            fallback_query = " ".join(meaningful_keywords)
            print(f"📝 Using title-based query: {fallback_query}")
            return fallback_query
        
        # Step 4: Ultimate fallback
        print(f"⚠️ Using generic fallback")
        return "business presentation professional"
    
    def _is_diagram_query(self, query: str) -> bool:
        """Check if query is looking for diagrams/illustrations"""
        diagram_keywords = [
            'diagram', 'architecture', 'flowchart', 'schema', 'chart',
            'graph', 'topology', 'illustration', 'structure', 'model',
            'uml', 'er diagram', 'network', 'system', 'design', 'workflow',
            'layer', 'kernel', 'process', 'memory'
        ]
        query_lower = query.lower()
        return any(keyword in query_lower for keyword in diagram_keywords)
    
    def _get_image_from_google(self, query: str):
        """
        Get image from Google Custom Search API
        Best for: Technical diagrams, educational content, specific illustrations
        """
        try:
            if not self.google_api_key or not self.google_cx:
                print("✗ Google API credentials not found")
                return None
            
            if not query or query.strip() == "":
                query = "technology diagram"
            
            print(f"🔍 Searching Google Images: '{query}'")
            
            url = "https://www.googleapis.com/customsearch/v1"
            params = {
                "key": self.google_api_key,
                "cx": self.google_cx,
                "q": query,
                "searchType": "image",
                "num": 5,
                "imgSize": "large",
                "imgType": "photo",
                "safe": "active",
                "fileType": "jpg,png"
            }
            
            response = requests.get(url, params=params, timeout=10)
            
            if response.status_code == 200:
                data = response.json()
                if data.get("items") and len(data["items"]) > 0:
                    # Try each result until one downloads successfully
                    for item in data["items"][:3]:
                        try:
                            image_url = item["link"]
                            img_response = requests.get(image_url, timeout=15, headers={
                                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                            })
                            if img_response.status_code == 200:
                                print(f"✓ Image downloaded from Google Images")
                                return BytesIO(img_response.content)
                        except:
                            continue
                    
                    print(f"⚠️ Found results but couldn't download images")
                else:
                    print(f"⚠️ No results found on Google for '{query}'")
            elif response.status_code == 429:
                print(f"⚠️ Google API rate limit reached")
            else:
                print(f"✗ Google API error: {response.status_code}")
            
            return None
                
        except Exception as e:
            print(f"✗ Google search error: {e}")
            return None
    
    def _get_image_from_pixabay(self, query: str):
        """
        Get image from Pixabay API
        Supports both photos and illustrations/vectors for diagrams
        """
        try:
            if not self.pixabay_api_key:
                print("✗ Pixabay API key not found")
                return None
            
            if not query or query.strip() == "":
                query = "business presentation"
            
            # Determine if we need illustrations or photos
            is_diagram = self._is_diagram_query(query)
            
            if is_diagram:
                print(f"🔍 Searching Pixabay for ILLUSTRATIONS: '{query}'")
                image_type = "vector,illustration"
                # Add editor's choice for better quality diagrams
                editors_choice = "true"
            else:
                print(f"🔍 Searching Pixabay for PHOTOS: '{query}'")
                image_type = "photo"
                editors_choice = "false"
            
            url = "https://pixabay.com/api/"
            params = {
                "key": self.pixabay_api_key,
                "q": query,
                "image_type": image_type,
                "orientation": "horizontal",
                "per_page": 5,  # Get more results to choose from
                "safesearch": "true",
                "min_width": 1200,
                "editors_choice": editors_choice
            }
            
            response = requests.get(url, params=params, timeout=10)
            
            if response.status_code == 200:
                data = response.json()
                if data.get("hits") and len(data["hits"]) > 0:
                    # For diagrams, prioritize vector graphics
                    best_hit = data["hits"][0]
                    if is_diagram and len(data["hits"]) > 1:
                        # Look for vector type first
                        for hit in data["hits"]:
                            if hit.get("type") == "vector":
                                best_hit = hit
                                break
                    
                    image_url = best_hit["largeImageURL"]
                    img_response = requests.get(image_url, timeout=15)
                    if img_response.status_code == 200:
                        img_type = best_hit.get("type", image_type)
                        print(f"✓ Image downloaded from Pixabay (type: {img_type})")
                        return BytesIO(img_response.content)
                else:
                    print(f"⚠️ No results found on Pixabay for '{query}'")
                    # If no illustration results, try with photo type for backup
                    if is_diagram and image_type != "photo":
                        print(f"🔄 Retrying with broader search...")
                        params["image_type"] = "all"
                        response = requests.get(url, params=params, timeout=10)
                        if response.status_code == 200:
                            data = response.json()
                            if data.get("hits") and len(data["hits"]) > 0:
                                image_url = data["hits"][0]["largeImageURL"]
                                img_response = requests.get(image_url, timeout=15)
                                if img_response.status_code == 200:
                                    print(f"✓ Image downloaded (broader search)")
                                    return BytesIO(img_response.content)
            else:
                print(f"✗ Pixabay API error: {response.status_code}")
            
            return None
                
        except Exception as e:
            print(f"✗ Pixabay error: {e}")
            return None
    
    def _get_image_from_unsplash(self, query: str):
        """
        Get image from Unsplash API (No API key needed - using Source API)
        Fallback for general photos when Pixabay doesn't have results
        """
        try:
            if not query or query.strip() == "":
                query = "technology"
            
            print(f"🔄 Trying Unsplash fallback: '{query}'")
            
            # Using Unsplash Source API (no auth required)
            clean_query = query.replace(' ', ',')
            image_url = f"https://source.unsplash.com/1600x900/?{clean_query}"
            
            response = requests.get(image_url, timeout=10, allow_redirects=True)
            
            if response.status_code == 200:
                print(f"✓ Image downloaded from Unsplash")
                return BytesIO(response.content)
            
            return None
                
        except Exception as e:
            print(f"✗ Unsplash error: {e}")
            return None
    
    def _get_fallback_image(self, query: str):
        """Fallback to Lorem Picsum with seed based on query"""
        try:
            seed_num = abs(hash(query)) % 1000
            image_url = f"https://picsum.photos/seed/{seed_num}/1600/900"
            
            print(f"🔄 Using Lorem Picsum fallback...")
            response = requests.get(image_url, timeout=10, allow_redirects=True)
            
            if response.status_code == 200:
                print(f"✓ Fallback image fetched")
                return BytesIO(response.content)
            
            return None
        except Exception as e:
            print(f"✗ Lorem Picsum failed: {e}")
            return None
    
    def _get_image(self, query: str):
        """
        Get image with intelligent fallback chain
        Priority:
        1. Google Custom Search (best for technical diagrams from web)
        2. Pixabay (illustrations for diagrams, photos for general content)
        3. Unsplash (high-quality photos, no API key needed)
        4. Lorem Picsum (ultimate fallback)
        """
        is_diagram = self._is_diagram_query(query)
        
        # For diagrams/technical content, prioritize Google Search
        if is_diagram:
            result = self._get_image_from_google(query)
            if result:
                return result
        
        # Try Pixabay
        result = self._get_image_from_pixabay(query)
        if result:
            return result
        
        # Try Google if not already tried
        if not is_diagram:
            result = self._get_image_from_google(query)
            if result:
                return result
        
        # Fallback to Unsplash for general photos
        result = self._get_image_from_unsplash(query)
        if result:
            return result
        
        # Fallback to Lorem Picsum
        result = self._get_fallback_image(query)
        if result:
            return result
        
        # Final fallback - generic image
        try:
            print(f"⚠️ Using generic fallback...")
            response = requests.get("https://picsum.photos/1600/900", timeout=10, allow_redirects=True)
            if response.status_code == 200:
                return BytesIO(response.content)
        except:
            pass
        
        return None
    
    def add_title_slide(self, title: str, subtitle: str = ""):
        """Create title slide with full background image and overlay text"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Get background image
        print(f"📸 Fetching title slide background...")
        bg_query = f"{title} professional background"
        bg_image = self._get_image(bg_query)
        
        if bg_image:
            # ✅ FIX: Wrap in try-except to handle image errors
            try:
                # Add full background image
                pic = slide.shapes.add_picture(
                    bg_image,
                    Inches(0),
                    Inches(0),
                    width=self.prs.slide_width,
                    height=self.prs.slide_height
                )
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)
                
                # Add dark overlay for text readability
                overlay = slide.shapes.add_shape(1, 0, 0, self.prs.slide_width, self.prs.slide_height)
                overlay.fill.solid()
                overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
                overlay.fill.transparency = 0.4
                overlay.line.fill.background()
                slide.shapes._spTree.remove(overlay._element)
                slide.shapes._spTree.insert(3, overlay._element)
            except Exception as e:
                print(f"✗ Failed to add background image: {e}")
                # Fallback to solid background
                bg = slide.shapes.add_shape(1, 0, 0, self.prs.slide_width, self.prs.slide_height)
                bg.fill.solid()
                bg.fill.fore_color.rgb = self.theme_colors["secondary"]
                bg.line.fill.background()
        else:
            # Fallback solid background
            bg = slide.shapes.add_shape(1, 0, 0, self.prs.slide_width, self.prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.theme_colors["secondary"]
            bg.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = RGBColor(240, 240, 240)
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(self, title: str, content: List[str], notes: str = "", image_query: str = None):
        """Create split-screen slide: image on left, content on right"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Build smart query based on content analysis
        smart_query = self._build_smart_query(title, content, image_query)
        
        print(f"📸 Fetching image for: {title}")
        print(f"   Using query: '{smart_query}'")
        content_image = self._get_image(smart_query)
        
        # ✅ FIX: Only add image if successfully fetched
        if content_image:
            try:
                # Add image on LEFT side (40% width)
                pic = slide.shapes.add_picture(
                    content_image,
                    Inches(0),
                    Inches(0),
                    width=Inches(4),
                    height=self.prs.slide_height
                )
                print(f"✓ Image added to left side")
            except Exception as e:
                print(f"✗ Failed to add image: {e}")
                # Continue without image - just use white background
        
        # Add WHITE background on RIGHT side (60% width)
        right_bg = slide.shapes.add_shape(
            1,
            Inches(4),
            Inches(0),
            Inches(6),
            self.prs.slide_height
        )
        right_bg.fill.solid()
        right_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        right_bg.line.fill.background()
        
        # Add title on right side
        title_box = slide.shapes.add_textbox(
            Inches(4.3),
            Inches(0.5),
            Inches(5.4),
            Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(33, 33, 33)
        title_para.alignment = PP_ALIGN.LEFT
        
        # Add divider line under title
        divider = slide.shapes.add_shape(
            1,
            Inches(4.3),
            Inches(1.6),
            Inches(5.4),
            Inches(0.02)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.theme_colors["accent"]
        divider.line.fill.background()
        
        # Add content on right side
        content_box = slide.shapes.add_textbox(
            Inches(4.3),
            Inches(2),
            Inches(5.2),
            Inches(5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        
        # Calculate font size
        total_chars = sum(len(point) for point in content)
        num_points = len(content)
        
        if num_points <= 3 and total_chars < 300:
            base_font_size = 20
        elif num_points <= 4 and total_chars < 450:
            base_font_size = 18
        else:
            base_font_size = 16
        
        # Add bullet points
        for i, point in enumerate(content):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = f"• {point}"
            p.font.size = Pt(base_font_size)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_before = Pt(12)
            p.space_after = Pt(8)
            p.line_spacing = 1.4
        
        # Add notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes[:500]
    
    def add_two_column_slide(self, title: str, left_content: List[str], 
                            right_content: List[str], notes: str = "", image_query: str = None):
        """Create two-column slide with image background on left"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Build smart query
        all_content = left_content + right_content
        smart_query = self._build_smart_query(title, all_content, image_query)
        
        print(f"📸 Fetching image for two-column slide")
        print(f"   Using query: '{smart_query}'")
        bg_image = self._get_image(smart_query)
        
        # ✅ FIX: Only add image if successfully fetched
        if bg_image:
            try:
                # Add image on left 40%
                pic = slide.shapes.add_picture(
                    bg_image,
                    Inches(0),
                    Inches(0),
                    width=Inches(4),
                    height=self.prs.slide_height
                )
            except Exception as e:
                print(f"✗ Failed to add image: {e}")
        
        # White background on right
        right_bg = slide.shapes.add_shape(1, Inches(4), 0, Inches(6), self.prs.slide_height)
        right_bg.fill.solid()
        right_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        right_bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(4.3), Inches(0.5), Inches(5.4), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(33, 33, 33)
        
        # Divider
        divider = slide.shapes.add_shape(1, Inches(4.3), Inches(1.4), Inches(5.4), Inches(0.02))
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.theme_colors["accent"]
        divider.line.fill.background()
        
        # Left column on right side
        left_box = slide.shapes.add_textbox(Inches(4.3), Inches(1.8), Inches(5.2), Inches(2.4))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        for i, point in enumerate(left_content):
            if i > 0:
                left_frame.add_paragraph()
            p = left_frame.paragraphs[i]
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_before = Pt(10)
            p.line_spacing = 1.3
        
        # Right column on right side
        right_box = slide.shapes.add_textbox(Inches(4.3), Inches(4.5), Inches(5.2), Inches(2.4))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        for i, point in enumerate(right_content):
            if i > 0:
                right_frame.add_paragraph()
            p = right_frame.paragraphs[i]
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_before = Pt(10)
            p.line_spacing = 1.3
        
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes[:500]
    
    def save(self, filepath: str):
        """Save presentation"""
        self.prs.save(filepath)


def create_presentation(presentation_data: Dict, output_path: str, template: str = "modern"):
    """
    Create presentation with split-screen layout (image left + content right)
    Enhanced with smart image detection using Pixabay API
    """
    print(f"\n{'='*60}")
    print(f"Creating presentation: {presentation_data['title']}")
    print(f"{'='*60}\n")
    
    generator = PPTXGenerator(template=template)
    
    # Title slide with full background
    print(f"[Title Slide] {presentation_data['title']}")
    generator.add_title_slide(
        presentation_data["title"],
        "Generated by AutoSlideX"
    )
    
    # Process each slide with split-screen layout
    for idx, slide_data in enumerate(presentation_data["slides"], 1):
        layout_type = slide_data.get("layout_type", "content")
        title = slide_data["title"]
        content = slide_data["content"]
        notes = slide_data.get("notes", "")
        image_query = slide_data.get("image_query")
        
        print(f"\n[Slide {idx}/{len(presentation_data['slides'])}] {title}")
        
        if layout_type == "title":
            generator.add_title_slide(title, content[0] if content else "")
        elif layout_type == "two_column":
            mid = len(content) // 2
            generator.add_two_column_slide(title, content[:mid], content[mid:], notes, image_query)
        else:
            generator.add_content_slide(title, content, notes, image_query)
    
    # Thank you slide
    print(f"\n[Thank You Slide]")
    generator.add_title_slide("Thank You", "Questions?")
    
    # Save
    generator.save(output_path)
    
    print(f"\n{'='*60}")
    print(f"✅ Presentation saved: {output_path}")
    print(f"📊 Total slides: {len(presentation_data['slides']) + 2}")
    print(f"🎨 Layout: Split-screen (image left + content right)")
    print(f"🧠 Smart image detection: ENABLED")
    print(f"🔍 Image sources: Google Custom Search (primary for diagrams) + Pixabay + Unsplash")
    print(f"{'='*60}\n")
    
    return output_path